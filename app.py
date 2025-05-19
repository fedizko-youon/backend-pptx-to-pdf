from fastapi import FastAPI, HTTPException, BackgroundTasks, UploadFile, File, Form
from fastapi.responses import FileResponse
from pptx import Presentation
from pptx.shapes.graphfrm import GraphicFrame
import tempfile
import os
import json

app = FastAPI()

# Substitui texto nos shapes (textos e tabelas)
def substituir_texto_em_shape(shape, substituicoes):
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                texto_original = run.text
                texto_limpo = texto_original.strip()
                for chave, valor in substituicoes.items():
                    if chave in texto_limpo:
                        run.text = texto_original.replace(chave, valor).strip()
    elif isinstance(shape, GraphicFrame) and shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        texto_original = run.text
                        texto_limpo = texto_original.strip()
                        for chave, valor in substituicoes.items():
                            if chave in texto_limpo:
                                run.text = texto_original.replace(chave, valor).strip()

# Aplica substituições no arquivo PPTX
def substituir_em_apresentacao(caminho_entrada, caminho_saida, substituicoes):
    prs = Presentation(caminho_entrada)
    for slide in prs.slides:
        for shape in slide.shapes:
            substituir_texto_em_shape(shape, substituicoes)
    prs.save(caminho_saida)

# Remove arquivos temporários
def remover_arquivos(*caminhos):
    for caminho in caminhos:
        if os.path.exists(caminho):
            os.remove(caminho)

# Endpoint principal
@app.post("/editar/")
async def editar_pptx_upload(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    substituicoes_json: str = Form(...)
):
    try:
        try:
            substituicoes = json.loads(substituicoes_json)
        except json.JSONDecodeError:
            raise HTTPException(status_code=400, detail="JSON de substituições inválido.")

        # Cria arquivos temporários (entrada e saída)
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_input:
            temp_input.write(await file.read())
            temp_input_path = temp_input.name

        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_output:
            temp_output_path = temp_output.name

        # Aplica substituições
        substituir_em_apresentacao(temp_input_path, temp_output_path, substituicoes)

        # Nome do arquivo final
        nome_cliente = substituicoes.get("nome_cliente", "Cliente")
        nome_cliente_sanitizado = "".join(c for c in nome_cliente if c.isalnum() or c in (" ", "_", "-")).strip()
        nome_final = f"Proposta Comercial {nome_cliente_sanitizado}.pptx"

        # Limpeza pós-resposta
        background_tasks.add_task(remover_arquivos, temp_input_path, temp_output_path)

        return FileResponse(
            temp_output_path,
            filename=nome_final,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Erro: {str(e)}")